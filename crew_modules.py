import os
import sys
import pandas as pd
from io import StringIO
from pptx import Presentation
from crewai import Agent, Task, Crew, Process, LLM
from crewai.tools import tool
from sklearn.ensemble import RandomForestClassifier
from sklearn.model_selection import train_test_split
from sklearn.metrics import accuracy_score

# ==============================================================================
# 1. DEFINE TOOLS (The Skills)
# ==============================================================================

class TeamTools:

    @tool("Save File")
    def save_file(filename: str, content: str):
        """Saves text content to a file. Useful for reports, plans, and tickets."""
        with open(filename, 'w') as f:
            f.write(content)
        return f"File saved: {filename}"

    @tool("Inspect CSV")
    def inspect_csv(file_path: str):
        """
        Reads the first 5 rows and data types of a CSV file.
        Useful for understanding the dataset structure before cleaning or analysis.
        Input: file_path (str)
        """
        try:
            df = pd.read_csv(file_path)
            buffer = StringIO()
            df.info(buf=buffer)
            info_str = buffer.getvalue()
            return f"First 5 rows:\n{df.head().to_string()}\n\nData Info:\n{info_str}\n\nDescription:\n{df.describe().to_string()}"
        except Exception as e:
            return f"Error reading CSV: {e}"

    @tool("Execute Python Code")
    def execute_python_code(code: str):
        """
        Executes the given Python code. 
        The code must be valid Python. 
        Standard output (print statements) is captured and returned.
        Variables created in the code are NOT persisted between calls unless saved to files.
        Useful for dynamic data cleaning, analysis, and plotting.
        Input: code (str)
        """
        old_stdout = sys.stdout
        redirected_output = sys.stdout = StringIO()
        try:
            # Pre-import common libraries for convenience
            import pandas as pd
            import numpy as np
            import sklearn
            
            exec(code, globals())
            sys.stdout = old_stdout
            return f"Code executed successfully.\nOutput:\n{redirected_output.getvalue()}"
        except Exception as e:
            sys.stdout = old_stdout
            return f"Error executing code: {e}"

    @tool("Generate PowerPoint")
    def create_pptx(title: str, summary: str, findings: str):
        """Creates a .pptx slide deck. Inputs: Title, Summary, Findings."""
        prs = Presentation()

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Project Update: Data Analysis"
        slide.placeholders[1].text = title

        # Slide 2: Summary
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Summary"
        slide.placeholders[1].text = summary

        # Slide 3: Model Findings
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Model Insights"
        slide.placeholders[1].text = findings

        prs.save('churn_presentation.pptx')
        return "Presentation saved as 'churn_presentation.pptx'"

# ==============================================================================
# 2. AGENT INITIALIZATION
# ==============================================================================

def init_llms(api_key):
    gemini_pro = LLM(
        model="gemini/gemini-2.5-pro",
        verbose=True,
        temperature=0.7,
        api_key=api_key
    )
    gemini_flash = LLM(
        model="gemini/gemini-2.5-flash",
        verbose=True,
        temperature=0.5,
        api_key=api_key
    )
    return gemini_pro, gemini_flash

def get_intake_crew(api_key, request):
    gemini_pro, gemini_flash = init_llms(api_key)

    intake = Agent(
        role='Project Intake Manager',
        goal='Validate project requirements.',
        backstory="You are a strict director who ensures projects have clear business value.",
        verbose=True,
        llm=gemini_pro
    )

    scrum = Agent(
        role='Scrum Master',
        goal='Manage project flow and Jira tickets.',
        backstory="You keep the team organized. You create ticket files for tracking.",
        tools=[TeamTools.save_file],
        verbose=True,
        llm=gemini_flash
    )

    pm = Agent(
        role='Senior Project Manager',
        goal='Create a comprehensive project plan.',
        backstory="You are an experienced PM who translates business needs into technical tasks. You ensure the Data Engineer and Data Scientist have clear instructions.",
        tools=[TeamTools.save_file],
        verbose=True,
        llm=gemini_pro
    )

    task_intake = Task(
        description=f"Review request: '{request}'. Approve if valid.",
        expected_output="Approval decision.",
        agent=intake
    )

    task_scrum = Task(
        description="Create a 'jira_ticket.txt' with project scope and timeline.",
        expected_output="Confirmation of file creation.",
        agent=scrum
    )

    task_pm = Task(
        description=f"Create a detailed 'project_plan.md' based on the request: '{request}'. Outline specific technical steps for data cleaning and analysis.",
        expected_output="Confirmation of file creation.",
        agent=pm
    )

    return Crew(
        agents=[intake, scrum, pm],
        tasks=[task_intake, task_scrum, task_pm],
        process=Process.sequential,
        verbose=True
    )

def get_data_crew(api_key, csv_path, request):
    gemini_pro, _ = init_llms(api_key)

    engineer = Agent(
        role='Data Engineer',
        goal='Prepare clean datasets based on requirements.',
        backstory="You are an expert Python programmer. You inspect data and write custom code to clean it.",
        tools=[TeamTools.inspect_csv, TeamTools.execute_python_code],
        verbose=True,
        llm=gemini_pro
    )

    task_eng = Task(
        description=f"""
        1. Inspect the dataset at '{csv_path}'.
        2. Write and execute Python code to clean the data based on this request: '{request}'.
        3. Ensure the cleaned data is saved to 'cleaned_data.csv'.
        4. Verify the file exists.
        """,
        expected_output="Confirmation that 'cleaned_data.csv' has been created and cleaned.",
        agent=engineer
    )

    return Crew(
        agents=[engineer],
        tasks=[task_eng],
        process=Process.sequential,
        verbose=True
    )

def get_analysis_crew(api_key, csv_path, request):
    gemini_pro, _ = init_llms(api_key)

    scientist = Agent(
        role='Senior Data Scientist',
        goal='Analyze data, perform statistical analyses and build predictive models if needed.',
        backstory="You are an expert Data Scientist. You write custom Python code (sklearn, pandas) to solve problems.",
        tools=[TeamTools.inspect_csv, TeamTools.execute_python_code],
        verbose=True,
        llm=gemini_pro
    )

    task_sci = Task(
        description=f"""
        1. Inspect the cleaned dataset at '{csv_path}'.
        2. Based on the request '{request}', write and execute Python code to perform the analysis or train a model.
        3. If training a model, report accuracy/metrics.
        4. If performing analysis, report key insights.
        """,
        expected_output="A summary of the analysis results, model performance, or key insights.",
        agent=scientist
    )

    return Crew(
        agents=[scientist],
        tasks=[task_sci],
        process=Process.sequential,
        verbose=True
    )

def get_reporting_crew(api_key, analysis_result):
    _, gemini_flash = init_llms(api_key)

    slide_maker = Agent(
        role='Presentation Designer',
        goal='Create visual slide decks.',
        backstory="You take technical results and turn them into PowerPoint files.",
        tools=[TeamTools.create_pptx],
        verbose=True,
        llm=gemini_flash
    )

    task_ppt = Task(
        description=f"Create a PowerPoint. Title: 'Project Analysis Results'. Use this summary: {analysis_result}",
        expected_output="Confirmation that .pptx is saved.",
        agent=slide_maker
    )

    return Crew(
        agents=[slide_maker],
        tasks=[task_ppt],
        process=Process.sequential,
        verbose=True
    )
