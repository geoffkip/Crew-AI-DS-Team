import os
import pandas as pd
from pptx import Presentation
from crewai import Agent, Task, Crew, Process, LLM
from crewai.tools import tool
from sklearn.ensemble import RandomForestClassifier
from sklearn.model_selection import train_test_split
from sklearn.metrics import accuracy_score

# ==============================================================================
# 0. CONFIGURATION & LLM SETUP
# ==============================================================================

# 1. The "Big Brain" - Gemini 2.5 Pro
# Best for: Complex reasoning, coding, writing plans.
gemini_pro = LLM(
    model="gemini/gemini-2.5-pro",
    verbose=True,
    temperature=0.7,
    api_key=os.getenv("GOOGLE_API_KEY")
)

# 2. The "Fast Brain" - Gemini 2.5 Flash
# Best for: Speed, formatting, summaries, simple tasks.
gemini_flash = LLM(
    model="gemini/gemini-2.5-flash",
    verbose=True,
    temperature=0.5,
    api_key=os.getenv("GOOGLE_API_KEY")
)

# ==============================================================================
# 1. DEFINE TOOLS (The Skills)
# ==============================================================================

class TeamTools:

    @tool("Save File")
    def save_file(filename: str, content: str):
        """Saves text content to a file. Useful for reports, plans, and tickets."""
        with open(filename, 'w') as f:
            f.write(content)
        return f"File saved: {filename}"

    @tool("Create & Clean Data")
    def create_data(filename: str):
        """Creates a dummy dataset for churn prediction, cleans it, and saves it."""
        # Generating synthetic data
        data = {
            'Age': [25, 30, 45, 35, 50, 23, 60, 48, 33, 29, 22, 55],
            'Annual_Spend': [1200, 3000, 5000, 2500, 6000, 1000, 7000, 5500, 2800, 2900, 800, 6200],
            'Support_Calls': [1, 0, 0, 1, 0, 5, 0, 0, 1, 0, 4, 0],
            'Churn': [1, 0, 0, 0, 0, 1, 0, 0, 1, 0, 1, 0] # 1 = Churn (Left), 0 = Stayed
        }
        df = pd.DataFrame(data)
        # Mock cleaning: ensure no negative numbers
        df = df[df['Annual_Spend'] > 0] 
        df.to_csv(filename, index=False)
        return f"Clean dataset created at: {filename}"

    @tool("Train Random Forest")
    def train_model(csv_file: str):
        """
        Trains a Random Forest Classifier on the CSV data.
        Target column must be 'Churn'. 
        Returns accuracy score and feature importance.
        """
        try:
            df = pd.read_csv(csv_file)
            X = df.drop('Churn', axis=1)
            y = df['Churn']

            X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.3, random_state=42)
            
            # Initialize and train Random Forest
            clf = RandomForestClassifier(n_estimators=100, random_state=42)
            clf.fit(X_train, y_train)
            
            # Predict
            preds = clf.predict(X_test)
            acc = accuracy_score(y_test, preds)
            
            # Get Feature Importance
            importances = dict(zip(X.columns, clf.feature_importances_))
            
            return f"Model Trained. Accuracy: {acc:.2f}. Feature Importance: {importances}"
        except Exception as e:
            return f"Error training model: {e}"

    @tool("Generate PowerPoint")
    def create_pptx(title: str, summary: str, findings: str):
        """Creates a .pptx slide deck. Inputs: Title, Summary, Findings."""
        prs = Presentation()

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Project Update: Churn Prediction"
        slide.placeholders[1].text = title

        # Slide 2: Summary
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Summary"
        slide.placeholders[1].text = summary

        # Slide 3: Model Findings
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Model Insights"
        slide.placeholders[1].text = findings

        prs.save('churn_presentation.pptx')
        return "Presentation saved as 'churn_presentation.pptx'"

# ==============================================================================
# 2. DEFINE AGENTS (The Team)
# ==============================================================================

# 1. Intake Manager (Strategy)
intake = Agent(
    role='Project Intake Manager',
    goal='Validate project requirements.',
    backstory="You are a strict director who ensures projects have clear business value.",
    verbose=True,
    llm=gemini_pro
)

# 2. Scrum Master (Operations)
scrum = Agent(
    role='Scrum Master',
    goal='Manage project flow and Jira tickets.',
    backstory="You keep the team organized. You create ticket files for tracking.",
    tools=[TeamTools.save_file],
    verbose=True,
    llm=gemini_flash
)

# 3. Data Engineer (Data Prep)
engineer = Agent(
    role='Data Engineer',
    goal='Prepare clean datasets.',
    backstory="You build robust pipelines. You create the CSV files for the team.",
    tools=[TeamTools.create_data],
    verbose=True,
    llm=gemini_pro
)

# 4. Data Scientist (Analysis)
scientist = Agent(
    role='Senior Data Scientist',
    goal='Build predictive models.',
    backstory="You are an expert in Scikit-Learn. You interpret model accuracy and feature importance.",
    tools=[TeamTools.train_model],
    verbose=True,
    llm=gemini_pro
)

# 5. Slide Maker (Reporting)
slide_maker = Agent(
    role='Presentation Designer',
    goal='Create visual slide decks.',
    backstory="You take technical results and turn them into PowerPoint files.",
    tools=[TeamTools.create_pptx],
    verbose=True,
    llm=gemini_flash
)

# ==============================================================================
# 3. DEFINE TASKS (The Workflow)
# ==============================================================================

# Step 1: Approve
task_intake = Task(
    description="Review request: 'Build a churn model to see why users leave'. Approve if valid.",
    expected_output="Approval decision.",
    agent=intake
)

# Step 2: Track
task_scrum = Task(
    description="Create a 'jira_ticket.txt' with project scope and timeline.",
    expected_output="Confirmation of file creation.",
    agent=scrum
)

# Step 3: Build Data
task_eng = Task(
    description="Create the dataset 'churn_data.csv' with columns: Age, Spend, Support_Calls, Churn.",
    expected_output="Confirmation that data is ready.",
    agent=engineer
)

# Step 4: Train Model
task_sci = Task(
    description="Train a Random Forest on 'churn_data.csv'. Report the Accuracy and which feature is most important.",
    expected_output="A summary of the model performance and key drivers.",
    agent=scientist
)

# Step 5: Make Slides
task_ppt = Task(
    description="Create a PowerPoint. Title: 'Q4 Churn Analysis'. Use the Scientist's summary for content.",
    expected_output="Confirmation that .pptx is saved.",
    agent=slide_maker,
    context=[task_sci] # Explicitly pass the scientist's output to the slide maker
)

# ==============================================================================
# 4. EXECUTE THE CREW
# ==============================================================================

churn_crew = Crew(
    agents=[intake, scrum, engineer, scientist, slide_maker],
    tasks=[task_intake, task_scrum, task_eng, task_sci, task_ppt],
    process=Process.sequential,
    verbose=True
)

print("### STARTING GEMINI 2.5 AGENT TEAM ###")
result = churn_crew.kickoff()

print("\n\n########################")
print("## FINAL OUTPUT ##")
print("########################\n")
print(result)